"""Export API endpoints for notebook export"""
from fastapi import APIRouter, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel
from typing import Optional, List, Any
from storage.notebook_store import notebook_store
from storage.source_store import source_store
from storage.highlights_store import highlights_store
from datetime import datetime

router = APIRouter()


class ChatHistoryItem(BaseModel):
    """Chat history item for export"""
    question: str
    answer: str
    citations: Optional[List[Any]] = None
    timestamp: Optional[str] = None


class ExportRequest(BaseModel):
    """Export request - matches frontend ExportOptions"""
    notebook_id: str
    format: str  # 'markdown', 'html', 'pdf', 'pptx'
    include_sources_content: Optional[bool] = False
    chat_history: Optional[List[ChatHistoryItem]] = None
    pptx_theme: Optional[str] = "light"  # light, dark, corporate, academic


# Built-in PPTX color themes
PPTX_THEMES = {
    "light": {
        "bg": (255, 255, 255),
        "title": (30, 30, 30),
        "body": (60, 60, 60),
        "accent": (59, 130, 246),   # blue-500
        "footer": (150, 150, 150),
    },
    "dark": {
        "bg": (30, 30, 40),
        "title": (240, 240, 245),
        "body": (200, 200, 210),
        "accent": (99, 179, 237),   # sky-400
        "footer": (100, 100, 120),
    },
    "corporate": {
        "bg": (245, 245, 250),
        "title": (15, 40, 80),
        "body": (40, 50, 70),
        "accent": (0, 82, 155),     # navy
        "footer": (130, 130, 140),
    },
    "academic": {
        "bg": (253, 251, 247),
        "title": (60, 40, 20),
        "body": (70, 55, 35),
        "accent": (139, 69, 19),    # saddle brown
        "footer": (140, 130, 120),
    },
}


class ExportFormat(BaseModel):
    """Export format definition"""
    id: str
    name: str
    extension: str
    description: str


class SlideData(BaseModel):
    """A single slide's content."""
    title: str
    bullets: List[str] = []
    slide_type: str = "content"  # title, content, sources, qa, thankyou, visual_overview
    mermaid_code: Optional[str] = None  # Mermaid diagram code for visual_overview slides
    speaker_notes: Optional[str] = None  # Presenter notes for this slide


class PptxPreviewRequest(BaseModel):
    """Request to generate a slide preview (JSON, not bytes)."""
    notebook_id: str
    pptx_theme: Optional[str] = "light"


class PptxReviseRequest(BaseModel):
    """Request to revise slides using a natural language prompt."""
    notebook_id: str
    slides: List[SlideData]
    revision_prompt: str
    pptx_theme: Optional[str] = "light"


class PptxDownloadRequest(BaseModel):
    """Request to download final slides as .pptx bytes."""
    notebook_id: str
    slides: List[SlideData]
    pptx_theme: Optional[str] = "light"


@router.get("/formats")
async def get_export_formats():
    """Get available export formats"""
    formats = [
        {
            "id": "markdown",
            "name": "Markdown",
            "extension": "md",
            "description": "Plain text with formatting, great for notes and documentation"
        },
        {
            "id": "html",
            "name": "HTML",
            "extension": "html",
            "description": "Web page format, viewable in any browser"
        },
        {
            "id": "pdf",
            "name": "PDF",
            "extension": "pdf",
            "description": "Portable document format, best for sharing and printing"
        },
        {
            "id": "pptx",
            "name": "PowerPoint",
            "extension": "pptx",
            "description": "Presentation slides, great for meetings and sharing insights"
        }
    ]
    return {"formats": formats}


@router.post("/notebook")
async def export_notebook(request: ExportRequest):
    """Export a notebook to the specified format"""
    notebook = await notebook_store.get(request.notebook_id)
    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook not found")
    
    sources = await source_store.list(request.notebook_id)
    
    # Get highlights for each source
    all_highlights = []
    for source in sources:
        highlights = await highlights_store.list(request.notebook_id, source["id"])
        for h in highlights:
            h["source_filename"] = source.get("filename", "Unknown")
        all_highlights.extend(highlights)
    
    if request.format == "markdown":
        content = _generate_markdown(notebook, sources, all_highlights, request.chat_history, request.include_sources_content)
        return Response(
            content=content,
            media_type="text/markdown",
            headers={"Content-Disposition": f"attachment; filename={notebook['title']}.md"}
        )
    
    elif request.format == "html":
        content = _generate_html(notebook, sources, all_highlights, request.chat_history, request.include_sources_content)
        return Response(
            content=content,
            media_type="text/html",
            headers={"Content-Disposition": f"attachment; filename={notebook['title']}.html"}
        )
    
    elif request.format == "pdf":
        # For PDF, we return HTML and let the frontend handle PDF generation
        # (using jsPDF as shown in the frontend code)
        content = _generate_html(notebook, sources, all_highlights, request.chat_history, request.include_sources_content)
        return Response(
            content=content,
            media_type="text/html",
            headers={"Content-Disposition": f"attachment; filename={notebook['title']}.html"}
        )
    
    elif request.format == "pptx":
        content = await _generate_smart_pptx(notebook, sources, all_highlights, request.chat_history, request.pptx_theme or "light")
        return Response(
            content=content,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={notebook['title']}.pptx"}
        )
    
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported format: {request.format}")


# =============================================================================
# PPTX Prompt-Based Revision Endpoints
# =============================================================================

@router.post("/pptx/preview")
async def pptx_preview(request: PptxPreviewRequest):
    """Generate AI slide content as JSON for frontend preview.
    
    Returns a list of slide objects the frontend can render as cards.
    The user can then revise them with natural language prompts.
    """
    notebook = await notebook_store.get(request.notebook_id)
    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook not found")
    
    slides = await _generate_slides_json(notebook)
    return {"slides": [s.dict() for s in slides], "theme": request.pptx_theme}


@router.post("/pptx/revise")
async def pptx_revise(request: PptxReviseRequest):
    """Revise existing slides using a natural language prompt.
    
    Takes the current slides + a revision instruction and returns updated slides.
    This is the NotebookLM-style prompt-based revision feature.
    """
    notebook = await notebook_store.get(request.notebook_id)
    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook not found")
    
    revised = await _revise_slides_with_prompt(request.slides, request.revision_prompt, notebook)
    return {"slides": [s.dict() for s in revised], "theme": request.pptx_theme}


class MermaidPreviewRequest(BaseModel):
    """Request to render a Mermaid diagram to PNG for preview."""
    mermaid_code: str


@router.post("/pptx/render-diagram")
async def render_diagram_preview(request: MermaidPreviewRequest):
    """Render Mermaid code to PNG image for the slide builder preview.
    
    Returns base64-encoded PNG or an error message.
    """
    import base64
    try:
        from services.mermaid_renderer import render_mermaid_to_png, is_available
        if not is_available():
            return {"success": False, "error": "Playwright not available"}
        
        png_bytes = await render_mermaid_to_png(request.mermaid_code, width=800, height=500, scale=1.5)
        if png_bytes:
            b64 = base64.b64encode(png_bytes).decode("ascii")
            return {"success": True, "image": f"data:image/png;base64,{b64}"}
        return {"success": False, "error": "Render returned empty"}
    except Exception as e:
        return {"success": False, "error": str(e)}


@router.post("/pptx/download")
async def pptx_download(request: PptxDownloadRequest):
    """Convert finalized slide JSON into a downloadable .pptx file."""
    notebook = await notebook_store.get(request.notebook_id)
    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook not found")
    
    content = await _render_slides_to_pptx(request.slides, notebook, request.pptx_theme or "light")
    return Response(
        content=content,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={notebook['title']}.pptx"}
    )


async def _generate_slides_json(notebook: dict) -> List[SlideData]:
    """Use context_builder + LLM to generate slide content as structured data."""
    import json
    slides = [
        SlideData(
            title=notebook["title"],
            bullets=[notebook.get("description") or f"Exported: {datetime.now().strftime('%Y-%m-%d')}"],
            slide_type="title",
        )
    ]
    
    try:
        from services.context_builder import context_builder
        from services.rag_engine import rag_engine
        
        built = await context_builder.build_context(
            notebook_id=notebook["id"],
            skill_id="briefing",
            topic=notebook.get("description") or notebook["title"],
        )
        
        if built.sources_used > 0:
            prompt = f"""You are a slide content generator. Given the source material below, create presentation slides.

Return ONLY valid JSON — an array of objects, each with:
  "title": slide title (short),
  "bullets": array of 3-4 concise bullet points,
  "speaker_notes": 2-3 sentences the presenter should say for this slide (include specific details, numbers, or context from the source material that supports the bullets)

Create these slides in order:
1. "Executive Summary" — 3-4 key takeaways from all sources
2-4. One slide per major theme/topic found (2-4 slides)

Source material:
{built.context[:12000]}

JSON array:"""
            
            raw = await rag_engine._call_ollama(
                "You output only valid JSON arrays. No markdown, no explanation.",
                prompt,
                num_predict=1500,
                temperature=0.35,
            )
            
            start = raw.find("[")
            end = raw.rfind("]") + 1
            if start >= 0 and end > start:
                ai_slides = json.loads(raw[start:end])
                for s in ai_slides[:6]:
                    slides.append(SlideData(
                        title=s.get("title", "Key Finding"),
                        bullets=s.get("bullets", [])[:5],
                        slide_type="content",
                        speaker_notes=s.get("speaker_notes"),
                    ))
                print(f"[PPTX Preview] AI generated {len(ai_slides)} content slides")
                
                # Generate a real Mermaid diagram for the visual overview slide
                theme_titles = [s.get("title", "") for s in ai_slides[:6] if s.get("title")]
                if len(theme_titles) >= 2:
                    mermaid_code = await _generate_overview_mermaid(
                        notebook["title"], theme_titles, built.context[:4000]
                    )
                    slides.append(SlideData(
                        title="Visual Overview",
                        bullets=theme_titles,
                        slide_type="visual_overview",
                        mermaid_code=mermaid_code,
                    ))
    except Exception as e:
        print(f"[PPTX Preview] AI generation failed: {e}")
        slides.append(SlideData(title="Key Points", bullets=["(AI generation failed — add content manually)"], slide_type="content"))
    
    slides.append(SlideData(title="Thank You", bullets=[f"Generated from {notebook['title']} by LocalBook"], slide_type="thankyou"))
    return slides


async def _generate_overview_mermaid(title: str, themes: list, context: str) -> Optional[str]:
    """Generate a Mermaid diagram summarizing notebook themes.
    
    Adapts diagram type to content:
    - Metric-heavy content → pie chart or xychart-beta bar
    - Comparison/pros-cons → flowchart with two columns
    - General themes → mindmap with sub-points
    
    Applies color intelligence via %%{init}%% theme overrides based on
    detected content domain (financial=blue, growth=green, etc.)
    """
    from services.visual_generator import detect_content_domain, DOMAIN_PALETTES
    from services.visual_analyzer import visual_analyzer
    
    try:
        from services.rag_engine import rag_engine
        
        # Detect content characteristics for diagram type selection
        all_text = " ".join(themes) + " " + context[:2000]
        metric_result = visual_analyzer._detect_metrics_fast(all_text)
        domain = detect_content_domain(all_text)
        palette = DOMAIN_PALETTES.get(domain, DOMAIN_PALETTES["neutral"])
        
        # Route to the right diagram type
        if metric_result["is_metric_heavy"] and len(metric_result["metrics"]) >= 3:
            diagram_type = metric_result["metric_type"]  # "distribution", "trend_chart", or "key_stats"
        else:
            diagram_type = "mindmap"
        
        print(f"[PPTX] Diagram routing: domain={domain}, metrics={len(metric_result['metrics'])}, type={diagram_type}")
        
        if diagram_type == "distribution":
            # Pie chart from percentage metrics
            metrics = metric_result["metrics"][:7]
            lines = ['pie showData', f'    title {title[:50]}']
            for m in metrics:
                safe_label = m["label"].replace('"', "'")[:30]
                lines.append(f'    "{safe_label}" : {m["value"]}')
            code = "\n".join(lines)
            
        elif diagram_type in ("trend_chart", "key_stats") and len(metric_result["metrics"]) >= 3:
            # Bar chart via xychart-beta
            metrics = metric_result["metrics"][:8]
            labels = [f'"{m["label"][:20]}"' for m in metrics]
            values = [str(m["value"]) for m in metrics]
            code = f"""xychart-beta
    title "{title[:50]}"
    x-axis [{", ".join(labels)}]
    bar [{", ".join(values)}]"""
            
        else:
            # Mindmap with LLM-generated sub-points
            theme_list = "\n".join(f"  - {t}" for t in themes[:6])
            prompt = f"""Create a Mermaid mindmap diagram for a presentation.

Root node: {title}
Branches (one per theme):
{theme_list}

Source context for sub-points:
{context[:3000]}

Rules:
- Use "mindmap" as the diagram type
- Root node uses (( )) syntax
- Each theme is a direct child of root
- Add 1-2 short sub-points under each theme from the source material
- Keep all labels under 40 characters
- No special characters that break Mermaid (no quotes, no parentheses in labels)
- Output ONLY the Mermaid code, nothing else

Example format:
mindmap
  root((Topic Name))
    Theme One
      Key detail A
      Key detail B
    Theme Two
      Key detail C"""

            raw = await rag_engine._call_ollama(
                "You output only valid Mermaid diagram code. No markdown fences, no explanation.",
                prompt,
                num_predict=800,
                temperature=0.3,
            )
            
            code = raw.strip()
            if code.startswith("```"):
                code = code.split("\n", 1)[1] if "\n" in code else code[3:]
            if code.endswith("```"):
                code = code[:-3]
            code = code.strip()
            
            if not code.lower().startswith("mindmap"):
                code = f"mindmap\n{code}"
        
        # Inject color intelligence via Mermaid init directive
        # This applies domain-appropriate colors to the rendered diagram
        accent = palette["accent"]
        primary_fill = palette["fills"][0]
        primary_stroke = palette["strokes"][0]
        secondary_fill = palette["fills"][2] if len(palette["fills"]) > 2 else palette["fills"][0]
        
        init_block = f"""%%{{init: {{'theme': 'base', 'themeVariables': {{'primaryColor': '{primary_fill}', 'primaryBorderColor': '{primary_stroke}', 'primaryTextColor': '{accent}', 'secondaryColor': '{secondary_fill}', 'lineColor': '{primary_stroke}'}}}}}}%%
"""
        code = init_block + code
        
        print(f"[PPTX] Generated {diagram_type} diagram ({len(code)} chars, domain={domain})")
        return code
        
    except Exception as e:
        print(f"[PPTX] Mermaid generation failed: {e}")
        branches = "\n".join(f"    {t}" for t in themes[:6])
        safe_title = title.replace("(", "").replace(")", "").replace('"', "")[:40]
        return f"mindmap\n  root(({safe_title}))\n{branches}"


async def _revise_slides_with_prompt(slides: List[SlideData], revision_prompt: str, notebook: dict) -> List[SlideData]:
    """Revise slide content using LLM with the user's natural language instruction.
    
    Key design decisions:
    - Includes notebook source context so the LLM can add/expand content from real data
    - Preserves mermaid_code on visual_overview slides unless explicitly changed
    - Regenerates diagram if visual_overview bullets change
    """
    import json
    
    current_json = json.dumps([
        {"title": s.title, "bullets": s.bullets, "slide_type": s.slide_type}
        for s in slides
    ], indent=2)
    
    # Fetch source context so revisions can reference actual notebook content
    source_context = ""
    try:
        from services.context_builder import context_builder
        built = await context_builder.build_context(
            notebook_id=notebook["id"],
            skill_id="briefing",
            topic=notebook.get("description") or notebook["title"],
        )
        if built.sources_used > 0:
            source_context = f"\n\nNOTEBOOK SOURCE MATERIAL (use this to add or expand content):\n{built.context[:6000]}"
    except Exception as e:
        print(f"[PPTX Revise] Could not fetch sources: {e}")
    
    prompt = f"""You are a presentation editor. Below are the current slides as JSON, the notebook's source material, and a revision instruction from the user.

Current slides:
{current_json}
{source_context}

User's revision instruction:
"{revision_prompt}"

Apply the revision and return the COMPLETE updated slide deck as a JSON array.
Each slide must have: "title", "bullets" (array of strings), "slide_type" (title/content/visual_overview/thankyou).
You may add, remove, reorder, or modify slides as the instruction requires.
When adding new content, USE the source material above — do not make things up.

Return ONLY valid JSON array:"""
    
    try:
        from services.rag_engine import rag_engine
        
        raw = await rag_engine._call_ollama(
            "You are a precise slide editor. Output only valid JSON arrays. No markdown, no explanation. Use source material for new content.",
            prompt,
            num_predict=2500,
            temperature=0.3,
        )
        
        start = raw.find("[")
        end = raw.rfind("]") + 1
        if start >= 0 and end > start:
            revised = json.loads(raw[start:end])
            
            # Build lookup of original mermaid_code by slide title
            original_mermaid = {
                s.title: s.mermaid_code for s in slides
                if s.mermaid_code and s.slide_type == "visual_overview"
            }
            
            result = []
            for s in revised:
                slide_type = s.get("slide_type", "content")
                new_slide = SlideData(
                    title=s.get("title", ""),
                    bullets=s.get("bullets", [])[:6],
                    slide_type=slide_type,
                )
                
                # Preserve or regenerate diagram for visual_overview slides
                if slide_type == "visual_overview":
                    old_code = original_mermaid.get(new_slide.title)
                    if old_code:
                        # Same title → keep existing diagram
                        new_slide.mermaid_code = old_code
                    else:
                        # New or changed visual_overview → regenerate diagram
                        try:
                            new_slide.mermaid_code = await _generate_overview_mermaid(
                                new_slide.title,
                                new_slide.bullets,
                                source_context[:4000],
                            )
                        except Exception:
                            pass  # Will fall back to colored boxes
                
                result.append(new_slide)
            
            print(f"[PPTX Revise] Applied revision: '{revision_prompt[:60]}' -> {len(result)} slides")
            return result
    except Exception as e:
        print(f"[PPTX Revise] Revision failed: {e}")
    
    return slides  # Return unchanged on failure


async def _render_slides_to_pptx(slides: List[SlideData], notebook: dict, theme_name: str) -> bytes:
    """Render structured slide data into a themed .pptx file.
    
    For visual_overview slides with mermaid_code, renders the diagram to PNG
    via Playwright and embeds it in the slide. Falls back to colored boxes
    if rendering is unavailable.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
    import io
    
    theme = PPTX_THEMES.get(theme_name, PPTX_THEMES["light"])
    bg_rgb = RgbColor(*theme["bg"])
    title_rgb = RgbColor(*theme["title"])
    body_rgb = RgbColor(*theme["body"])
    accent_rgb = RgbColor(*theme["accent"])
    footer_rgb = RgbColor(*theme["footer"])
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb
        
        if slide_data.slide_type in ("title", "thankyou"):
            box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
            p = box.text_frame.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = title_rgb
            p.alignment = PP_ALIGN.CENTER
            if slide_data.bullets:
                sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
                sp = sub.text_frame.paragraphs[0]
                sp.text = slide_data.bullets[0]
                sp.font.size = Pt(22 if slide_data.slide_type == "title" else 18)
                sp.font.color.rgb = accent_rgb
                sp.alignment = PP_ALIGN.CENTER
        
        elif slide_data.slide_type == "visual_overview":
            # Title
            tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
            tp = tbox.text_frame.paragraphs[0]
            tp.text = slide_data.title
            tp.font.size = Pt(32)
            tp.font.bold = True
            tp.font.color.rgb = title_rgb
            
            # PRIMARY: Render Mermaid diagram to PNG and embed as image
            png_embedded = False
            if slide_data.mermaid_code:
                try:
                    from services.mermaid_renderer import render_mermaid_to_png
                    png_bytes = await render_mermaid_to_png(slide_data.mermaid_code)
                    if png_bytes and len(png_bytes) > 100:
                        img_stream = io.BytesIO(png_bytes)
                        # Center the image on the slide below the title
                        # Calculate dimensions to fit within slide bounds
                        img_w = Inches(10)
                        img_h = Inches(5.2)
                        img_x = Inches(1.667)  # centered: (13.333 - 10) / 2
                        img_y = Inches(1.3)
                        slide.shapes.add_picture(img_stream, img_x, img_y, img_w, img_h)
                        png_embedded = True
                        print(f"[PPTX] Embedded {len(png_bytes)} byte PNG diagram in slide")
                except Exception as e:
                    print(f"[PPTX] PNG embed failed, falling back to boxes: {e}")
            
            # FALLBACK: Colored boxes if PNG rendering unavailable
            if not png_embedded:
                from services.visual_generator import detect_content_domain, DOMAIN_PALETTES
                
                all_text = " ".join(slide_data.bullets)
                domain = detect_content_domain(all_text)
                palette = DOMAIN_PALETTES.get(domain, DOMAIN_PALETTES["neutral"])
                
                items = slide_data.bullets[:6]
                cols = 3 if len(items) > 4 else 2
                box_w = 3.8
                box_h = 2.0
                gap_x = 0.3
                gap_y = 0.25
                start_x = (13.333 - (cols * box_w + (cols - 1) * gap_x)) / 2
                start_y = 1.5
                
                for idx, item in enumerate(items):
                    row = idx // cols
                    col = idx % cols
                    x = start_x + col * (box_w + gap_x)
                    y = start_y + row * (box_h + gap_y)
                    
                    shape = slide.shapes.add_shape(
                        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
                        Inches(x), Inches(y), Inches(box_w), Inches(box_h)
                    )
                    shape_fill = shape.fill
                    shape_fill.solid()
                    fill_idx = idx % len(palette["fills"])
                    hex_c = palette["fills"][fill_idx].lstrip("#")
                    shape_fill.fore_color.rgb = RgbColor(int(hex_c[:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16))
                    stroke_hex = palette["strokes"][fill_idx].lstrip("#")
                    shape.line.color.rgb = RgbColor(int(stroke_hex[:2], 16), int(stroke_hex[2:4], 16), int(stroke_hex[4:6], 16))
                    shape.line.width = Pt(1.5)
                    
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = item
                    p.font.size = Pt(14)
                    p.font.bold = True
                    stroke_rgb = RgbColor(int(stroke_hex[:2], 16), int(stroke_hex[2:4], 16), int(stroke_hex[4:6], 16))
                    p.font.color.rgb = stroke_rgb
                    p.alignment = PP_ALIGN.CENTER
        
        else:
            # Content slide
            tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
            tp = tbox.text_frame.paragraphs[0]
            tp.text = slide_data.title
            tp.font.size = Pt(36)
            tp.font.bold = True
            tp.font.color.rgb = title_rgb
            
            if slide_data.bullets:
                cbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
                tf = cbox.text_frame
                tf.word_wrap = True
                for i, bullet in enumerate(slide_data.bullets):
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.text = f"\u2022 {bullet}"
                    para.font.size = Pt(18)
                    para.font.color.rgb = body_rgb
                    para.space_after = Pt(10)
        
        # Speaker notes
        if slide_data.speaker_notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data.speaker_notes
        
        # Footer
        fbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4))
        fp = fbox.text_frame.paragraphs[0]
        fp.text = "Generated by LocalBook"
        fp.font.size = Pt(10)
        fp.font.color.rgb = footer_rgb
        fp.alignment = PP_ALIGN.RIGHT
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def _generate_markdown(notebook: dict, sources: list, highlights: list, chat_history: list, include_content: bool) -> str:
    """Generate markdown export"""
    lines = []
    
    # Title
    lines.append(f"# {notebook['title']}")
    lines.append("")
    
    if notebook.get('description'):
        lines.append(f"_{notebook['description']}_")
        lines.append("")
    
    lines.append(f"Exported: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    lines.append("")
    
    # Sources section
    lines.append("## Sources")
    lines.append("")
    
    if sources:
        for i, source in enumerate(sources, 1):
            lines.append(f"### {i}. {source.get('filename', 'Unknown')}")
            lines.append(f"- **Format:** {source.get('format', 'unknown').upper()}")
            lines.append(f"- **Chunks:** {source.get('chunks', 0)}")
            lines.append(f"- **Characters:** {source.get('characters', 0):,}")
            lines.append(f"- **Status:** {source.get('status', 'unknown')}")
            
            if source.get('url'):
                lines.append(f"- **URL:** {source['url']}")
            
            lines.append("")
            
            if include_content and source.get('content'):
                lines.append("#### Content")
                lines.append("```")
                lines.append(source['content'][:5000])  # Limit content length
                if len(source.get('content', '')) > 5000:
                    lines.append("... (truncated)")
                lines.append("```")
                lines.append("")
    else:
        lines.append("_No sources in this notebook_")
        lines.append("")
    
    # Highlights section
    if highlights:
        lines.append("## Highlights & Annotations")
        lines.append("")
        
        for h in highlights:
            lines.append(f"### From: {h.get('source_filename', 'Unknown')}")
            lines.append(f"> {h.get('highlighted_text', '')}")
            if h.get('annotation'):
                lines.append(f"**Note:** {h['annotation']}")
            lines.append("")
    
    # Chat history section
    if chat_history:
        lines.append("## Q&A History")
        lines.append("")
        
        for i, exchange in enumerate(chat_history, 1):
            lines.append(f"### Q{i}: {exchange.question}")
            if exchange.timestamp:
                lines.append(f"_Asked: {exchange.timestamp}_")
            lines.append("")
            lines.append(f"**Answer:** {exchange.answer}")
            lines.append("")
            
            if exchange.citations:
                lines.append("**Citations:**")
                for citation in exchange.citations:
                    lines.append(f"- [{citation.get('number', '?')}] {citation.get('filename', 'Unknown')}: {citation.get('snippet', '')[:100]}...")
                lines.append("")
    
    # Footer
    lines.append("---")
    lines.append("_Generated by LocalBook_")
    
    return "\n".join(lines)


def _generate_html(notebook: dict, sources: list, highlights: list, chat_history: list, include_content: bool) -> str:
    """Generate HTML export"""
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{notebook['title']} - LocalBook Export</title>
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, sans-serif;
            max-width: 800px;
            margin: 0 auto;
            padding: 2rem;
            line-height: 1.6;
            color: #333;
        }}
        h1 {{ color: #1a1a1a; border-bottom: 2px solid #3b82f6; padding-bottom: 0.5rem; }}
        h2 {{ color: #374151; margin-top: 2rem; }}
        h3 {{ color: #4b5563; }}
        .meta {{ color: #6b7280; font-size: 0.9rem; }}
        .source {{ background: #f9fafb; padding: 1rem; border-radius: 8px; margin: 1rem 0; }}
        .highlight {{ background: #fef3c7; padding: 1rem; border-left: 4px solid #f59e0b; margin: 1rem 0; }}
        .qa {{ background: #eff6ff; padding: 1rem; border-radius: 8px; margin: 1rem 0; }}
        .question {{ font-weight: bold; color: #1e40af; }}
        .answer {{ margin-top: 0.5rem; }}
        .citation {{ font-size: 0.85rem; color: #6b7280; margin-top: 0.5rem; }}
        blockquote {{ border-left: 4px solid #d1d5db; padding-left: 1rem; margin: 1rem 0; color: #4b5563; }}
        code {{ background: #f3f4f6; padding: 0.2rem 0.4rem; border-radius: 4px; }}
        pre {{ background: #1f2937; color: #f9fafb; padding: 1rem; border-radius: 8px; overflow-x: auto; }}
        .footer {{ margin-top: 3rem; padding-top: 1rem; border-top: 1px solid #e5e7eb; color: #9ca3af; font-size: 0.85rem; }}
    </style>
</head>
<body>
    <h1>{notebook['title']}</h1>
"""
    
    if notebook.get('description'):
        html += f"    <p class='meta'><em>{notebook['description']}</em></p>\n"
    
    html += f"    <p class='meta'>Exported: {datetime.now().strftime('%Y-%m-%d %H:%M')}</p>\n"
    
    # Sources
    html += "    <h2>📚 Sources</h2>\n"
    
    if sources:
        for i, source in enumerate(sources, 1):
            html += f"""    <div class='source'>
        <h3>{i}. {source.get('filename', 'Unknown')}</h3>
        <p><strong>Format:</strong> {source.get('format', 'unknown').upper()} | 
           <strong>Chunks:</strong> {source.get('chunks', 0)} | 
           <strong>Characters:</strong> {source.get('characters', 0):,}</p>
"""
            if source.get('url'):
                html += f"        <p><strong>URL:</strong> <a href='{source['url']}'>{source['url']}</a></p>\n"
            html += "    </div>\n"
    else:
        html += "    <p><em>No sources in this notebook</em></p>\n"
    
    # Highlights
    if highlights:
        html += "    <h2>🖍️ Highlights & Annotations</h2>\n"
        for h in highlights:
            html += f"""    <div class='highlight'>
        <p><strong>From:</strong> {h.get('source_filename', 'Unknown')}</p>
        <blockquote>{h.get('highlighted_text', '')}</blockquote>
"""
            if h.get('annotation'):
                html += f"        <p><strong>Note:</strong> {h['annotation']}</p>\n"
            html += "    </div>\n"
    
    # Chat history
    if chat_history:
        html += "    <h2>💬 Q&A History</h2>\n"
        for i, exchange in enumerate(chat_history, 1):
            html += f"""    <div class='qa'>
        <p class='question'>Q{i}: {exchange.question}</p>
"""
            if exchange.timestamp:
                html += f"        <p class='meta'>Asked: {exchange.timestamp}</p>\n"
            html += f"        <div class='answer'>{exchange.answer}</div>\n"
            
            if exchange.citations:
                html += "        <div class='citation'><strong>Citations:</strong><ul>\n"
                for citation in exchange.citations:
                    html += f"            <li>[{citation.get('number', '?')}] {citation.get('filename', 'Unknown')}</li>\n"
                html += "        </ul></div>\n"
            html += "    </div>\n"
    
    html += """    <div class='footer'>
        <p>Generated by LocalBook</p>
    </div>
</body>
</html>"""
    
    return html


async def _generate_smart_pptx(notebook: dict, sources: list, highlights: list, chat_history: list, theme_name: str = "light") -> bytes:
    """Generate AI-powered PowerPoint presentation.

    Builds SlideData objects and delegates to the unified _render_slides_to_pptx
    renderer, which provides all features: speaker notes, visual overview with
    Mermaid diagrams, color intelligence, and themed styling.

    Slide types:
      1. Title
      2-N. AI-generated content slides (with speaker notes)
      N+1. Visual overview (Mermaid diagram)
      N+2. Sources consulted
      N+3. Q&A highlights (if chat_history provided)
      Last. Thank You
    """
    # Build slide deck as structured data, then render through unified pipeline
    slides = await _generate_slides_json(notebook)
    
    # Insert sources slide before Thank You
    if sources:
        source_bullets = [
            f'{s.get("filename", "Unknown")} ({s.get("format", "").upper()})'
            for s in sources[:10]
        ]
        source_slide = SlideData(
            title="Sources Consulted",
            bullets=source_bullets,
            slide_type="sources",
        )
        # Insert before last slide (Thank You)
        if slides and slides[-1].slide_type == "thankyou":
            slides.insert(-1, source_slide)
        else:
            slides.append(source_slide)
    
    # Insert Q&A highlights before Thank You
    if chat_history:
        thankyou_idx = len(slides) - 1 if slides and slides[-1].slide_type == "thankyou" else len(slides)
        for i, exchange in enumerate(chat_history[:3], 1):
            q_text = exchange.question[:120] + ("..." if len(exchange.question) > 120 else "")
            answer = exchange.answer[:600]
            sentences = [s.strip() for s in answer.split(". ") if s.strip()][:5]
            qa_slide = SlideData(
                title=f"Q{i}: {q_text}",
                bullets=sentences,
                slide_type="qa",
                speaker_notes=f"This question was asked during the research session. Full answer: {exchange.answer[:300]}",
            )
            slides.insert(thankyou_idx, qa_slide)
            thankyou_idx += 1
    
    # Fallback: if no AI slides were generated, add highlights
    content_slides = [s for s in slides if s.slide_type == "content"]
    if not content_slides and highlights:
        highlight_slide = SlideData(
            title="Key Highlights",
            bullets=[
                f'"{h.get("highlighted_text", "")[:120]}..." — {h.get("source_filename", "")}'
                for h in highlights[:5]
            ],
            slide_type="content",
        )
        # Insert after title slide
        slides.insert(1, highlight_slide)
    
    return await _render_slides_to_pptx(slides, notebook, theme_name)
