"""Generate deterministic test files (PDF, PPTX, DOCX) for the LLM Evaluator.

Run once to create the test content:
    python -m evaluator.test_content.generate_test_files

Uses reportlab for PDF, python-pptx for PPTX, python-docx for DOCX.
Content is sourced from LocalBook's architecture docs to ensure relevance.
"""

import os
import sys
from pathlib import Path

OUTPUT_DIR = Path(__file__).parent


# ─── Content ─────────────────────────────────────────────────────────────────

PDF_CONTENT = """LocalBook System Architecture — Test Document

1. Overview

LocalBook is a privacy-first AI research notebook built with Tauri, FastAPI, and Ollama.
All processing happens locally on the user's machine using Apple Silicon GPUs.
No data is sent to external servers — the entire AI pipeline runs on-device.

The system architecture consists of five layers:
- Presentation Layer: Tauri webview with real-time UI
- API Layer: FastAPI backend with REST endpoints and WebSocket connections
- Intelligence Layer: RAG engine, context builder, intent classifier, and content generators
- Storage Layer: LanceDB for vector embeddings, JSON files for metadata
- Model Layer: Ollama-managed LLMs including reasoning, embedding, and vision models

2. RAG Pipeline

The Retrieval-Augmented Generation pipeline is the core of LocalBook's AI capabilities.
When a user asks a question, the system:
1. Embeds the query using the Snowflake Arctic Embed 2 model (1024 dimensions)
2. Searches the LanceDB vector store for the top-k most relevant chunks
3. Reranks results using cross-encoder scoring
4. Builds an adaptive context window based on available token budget
5. Generates an answer using the main reasoning model (Gemma 4 E4B by default)
6. Extracts and formats citations from the source material

The embedding model produces 1024-dimensional dense vectors. Each document is chunked
into passages of approximately 512 tokens with 50-token overlap. The chunking strategy
uses semantic boundaries (paragraphs, headers) rather than fixed character counts.

3. Model Roles

LocalBook uses a multi-model architecture with specialized roles:

Main Model (gemma4:e4b): Handles complex reasoning, document generation, native image
understanding, and streaming chat responses. Runs with a context window of up to 16384
tokens. Typical throughput: 15-25 tokens/second on M4 16GB. Apache-2.0 license.

Fast Model (phi4-mini): Used for quick tasks like intent classification, follow-up
questions, and simple lookups. 3.8B parameters with MIT license.

Embedding Model (snowflake-arctic-embed2): Generates vector embeddings for semantic
search. 568M parameters, produces 1024-dimensional vectors. Apache-2.0 license.

Vision: Gemma 4 E4B is natively multimodal, so it also analyzes images, charts, and
diagrams during document ingestion — no separate vision model is required by default.

4. Performance Benchmarks

On an Apple M4 chip with 16GB RAM:
- Gemma 4 E4B: 15-25 tokens/sec average, ~1s time-to-first-token
- Phi-4 Mini: 35 tokens/sec average, 0.6s time-to-first-token
- Snowflake Arctic Embed 2: 120 embeddings/sec for 512-token passages
- Gemma 4 vision: 3-5 seconds per image description
"""

PPTX_SLIDES = [
    {
        "title": "RAG Systems Overview",
        "content": "Retrieval-Augmented Generation combines information retrieval with language model generation.\n"
                   "Key components: Document store, Embedding model, Vector database, Language model, Context builder.\n"
                   "Benefits: Reduced hallucination, grounded answers, source citations, updatable knowledge."
    },
    {
        "title": "Embedding Models Compared",
        "content": "Snowflake Arctic Embed 2: 1024 dims, 568M params, fast throughput, Apache-2.0\n"
                   "Nomic Embed Text: 768 dims, 137M params, very small footprint, Apache-2.0\n"
                   "mxbai Embed Large: 1024 dims, 335M params, strong MTEB scores, Apache-2.0\n"
                   "All-MiniLM: 384 dims, 33M params, fastest but lower quality"
    },
    {
        "title": "Vector Search Architecture",
        "content": "Documents are chunked into ~512 token passages with semantic boundaries.\n"
                   "Each chunk is embedded into a dense vector using the embedding model.\n"
                   "Vectors are stored in LanceDB, an embedded vector database.\n"
                   "Queries are embedded and compared using cosine similarity.\n"
                   "Top-k results are reranked using cross-encoder scoring."
    },
    {
        "title": "Context Window Management",
        "content": "The adaptive context builder manages the token budget for each query.\n"
                   "7B models typically use 4096-8192 token contexts for optimal speed.\n"
                   "Retrieval chunks are prioritized by relevance score.\n"
                   "Source metadata (filename, page) is included for citation tracking.\n"
                   "Map-reduce overview provides global context from all sources."
    },
    {
        "title": "Performance on Apple Silicon",
        "content": "Apple M4 16GB: Supports 7B main + 3.8B fast model concurrently.\n"
                   "Metal GPU acceleration provides ~2x speedup over CPU inference.\n"
                   "Memory-mapped model loading enables fast cold starts.\n"
                   "Typical throughput: 15-25 tok/s (7B), 30-40 tok/s (3B).\n"
                   "TTFT: 1-2s for 7B models, <1s for 3B models."
    },
]

DOCX_CONTENT = """Embedding Model Evaluation Report

Executive Summary

This report evaluates embedding models for use in local RAG (Retrieval-Augmented Generation) systems deployed on consumer Apple Silicon hardware. The evaluation focuses on three key criteria: retrieval quality, throughput performance, and resource efficiency.

Methodology

We tested four embedding models across a corpus of 500 documents (research papers, web articles, and technical documentation). Each model was evaluated on:

1. MTEB benchmark scores for information retrieval tasks
2. Embedding throughput (passages per second) on Apple M4 16GB
3. Memory footprint during inference
4. Semantic discrimination accuracy on a custom test set

Results

Snowflake Arctic Embed 2 (1024 dimensions, 568M parameters):
- MTEB IR score: 0.72 (highest in class for its size)
- Throughput: 120 embeddings/sec on M4 16GB
- Memory: ~1.2 GB during inference
- Semantic discrimination: 94% accuracy on custom test
- License: Apache-2.0

Nomic Embed Text (768 dimensions, 137M parameters):
- MTEB IR score: 0.65
- Throughput: 200 embeddings/sec on M4 16GB
- Memory: ~0.3 GB during inference
- Semantic discrimination: 88% accuracy
- License: Apache-2.0

Recommendation

For LocalBook's use case (100-5000 documents per notebook, M4 16GB target hardware), Snowflake Arctic Embed 2 provides the best balance of quality and efficiency. The 1024-dimensional vectors offer superior retrieval accuracy while maintaining excellent throughput. The model's Apache-2.0 license ensures no deployment restrictions.

For users with limited RAM (8GB), Nomic Embed Text is a viable alternative with its smaller footprint, though retrieval quality is approximately 10% lower on our evaluation metrics.
"""


def generate_pdf():
    """Generate test_document.pdf using reportlab."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
        from reportlab.lib.units import inch
        from reportlab.lib.colors import HexColor
        from reportlab.graphics.shapes import Drawing, Rect, String
        from reportlab.graphics.charts.barcharts import VerticalBarChart
        from reportlab.graphics import renderPDF
    except ImportError:
        print("ERROR: reportlab not installed. Run: pip install reportlab")
        return False

    pdf_path = OUTPUT_DIR / "test_document.pdf"
    doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle('CustomTitle', parent=styles['Title'], fontSize=16, spaceAfter=12)
    heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading2'], fontSize=13, spaceAfter=8)
    body_style = ParagraphStyle('CustomBody', parent=styles['Normal'], fontSize=10, leading=14, spaceAfter=6)

    story = []

    # Parse content into sections
    lines = PDF_CONTENT.strip().split("\n")
    for line in lines:
        line = line.strip()
        if not line:
            story.append(Spacer(1, 6))
        elif line.startswith("LocalBook System Architecture"):
            story.append(Paragraph(line, title_style))
        elif line[0].isdigit() and ". " in line[:4]:
            story.append(Spacer(1, 12))
            story.append(Paragraph(line, heading_style))
        elif line.startswith("- "):
            story.append(Paragraph(f"• {line[2:]}", body_style))
        else:
            story.append(Paragraph(line, body_style))

    # Add a bar chart (for vision model testing)
    story.append(Spacer(1, 24))
    story.append(Paragraph("Figure 1: Model Throughput Comparison (tokens/sec)", heading_style))

    drawing = Drawing(400, 200)
    chart = VerticalBarChart()
    chart.x = 50
    chart.y = 30
    chart.width = 300
    chart.height = 140
    chart.data = [[18, 35, 120]]
    chart.categoryAxis.categoryNames = ['Gemma 4 E4B', 'Phi-4 Mini', 'Arctic Embed']
    chart.valueAxis.valueMin = 0
    chart.valueAxis.valueMax = 150
    chart.bars[0].fillColor = HexColor('#4A90D9')
    drawing.add(chart)
    story.append(drawing)

    doc.build(story)
    print(f"[TEST-CONTENT] Generated {pdf_path} ({pdf_path.stat().st_size} bytes)")
    return True


def generate_pptx():
    """Generate test_presentation.pptx using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx")
        return False

    pptx_path = OUTPUT_DIR / "test_presentation.pptx"
    prs = Presentation()

    for slide_data in PPTX_SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
        slide.shapes.title.text = slide_data["title"]

        # Add content to the text frame
        tf = slide.placeholders[1].text_frame
        tf.text = ""
        for i, line in enumerate(slide_data["content"].split("\n")):
            if i == 0:
                tf.text = line.strip()
            else:
                p = tf.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(14)

    prs.save(str(pptx_path))
    print(f"[TEST-CONTENT] Generated {pptx_path} ({pptx_path.stat().st_size} bytes)")
    return True


def generate_docx():
    """Generate test_report.docx using python-docx."""
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        print("ERROR: python-docx not installed. Run: pip install python-docx")
        return False

    docx_path = OUTPUT_DIR / "test_report.docx"
    doc = Document()

    lines = DOCX_CONTENT.strip().split("\n")
    for line in lines:
        line = line.strip()
        if not line:
            continue
        elif line == "Embedding Model Evaluation Report":
            doc.add_heading(line, level=0)
        elif line in ("Executive Summary", "Methodology", "Results", "Recommendation"):
            doc.add_heading(line, level=1)
        elif line.startswith(("Snowflake Arctic", "Nomic Embed")):
            doc.add_heading(line, level=2)
        elif line[0].isdigit() and ". " in line[:3]:
            doc.add_paragraph(line, style="List Number")
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)

    doc.save(str(docx_path))
    print(f"[TEST-CONTENT] Generated {docx_path} ({docx_path.stat().st_size} bytes)")
    return True


def main():
    """Generate all test files."""
    print("[TEST-CONTENT] Generating test files for LLM Evaluator...")

    results = {
        "PDF": generate_pdf(),
        "PPTX": generate_pptx(),
        "DOCX": generate_docx(),
    }

    success = sum(1 for v in results.values() if v)
    failed = sum(1 for v in results.values() if not v)
    print(f"\n[TEST-CONTENT] Done: {success} generated, {failed} failed")

    if failed:
        print("[TEST-CONTENT] Install missing dependencies:")
        print("  pip install reportlab python-pptx python-docx")
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
