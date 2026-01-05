"""Document processing service

All uploaded content is treated as source information for RAG retrieval.
The goal is to ensure every file type is accurately classified, stored, and
optimized for semantic search so questions can be answered from the content.
"""
import io
from typing import Dict
from pathlib import Path
from storage.source_store import source_store
from services.rag_engine import rag_engine

class DocumentProcessor:
    """Process and ingest documents"""

    async def process(self, content: bytes, filename: str, notebook_id: str) -> Dict:
        """Process a document and add to RAG system.
        
        Every uploaded file is treated as source information that should be
        retrievable when answering questions. This method:
        1. Extracts text optimized for semantic search (not just raw extraction)
        2. Creates a source record with metadata
        3. Ingests into the RAG vector store
        4. Extracts concepts for the knowledge graph
        """
        # Detect file type with magic byte fallback for unknown extensions
        file_format = self._get_file_type(filename, content)
        print(f"[DocProcessor] Processing {filename} (type: {file_format}, size: {len(content)} bytes)")

        # Extract text based on file type - optimized for semantic search
        text = await self._extract_text(content, filename)
        
        if not text or not text.strip():
            raise ValueError(f"No text content could be extracted from {filename}")
        
        print(f"[DocProcessor] Extracted {len(text)} characters from {filename}")

        # Create source record
        source = await source_store.create(
            notebook_id=notebook_id,
            filename=filename,
            metadata={
                "type": file_format,
                "format": file_format,  # Frontend expects 'format'
                "size": len(content),
                "chunks": 0,
                "characters": 0,
                "status": "processing"
            }
        )

        # Ingest into RAG system
        try:
            result = await rag_engine.ingest_document(
                notebook_id=notebook_id,
                source_id=source["id"],
                text=text,
                filename=filename,
                source_type=file_format
            )

            # Update source with processing results
            chunks = result.get("chunks", 0)
            characters = result.get("characters", len(text))

            # Save updated source back to store with content
            await source_store.update(notebook_id, source["id"], {
                "chunks": chunks,
                "characters": characters,
                "status": "completed",
                "content": text  # Save full text for viewing
            })

            # Update local copy for return
            source["chunks"] = chunks
            source["characters"] = characters
            source["status"] = "completed"

            return {
                "source_id": source["id"],
                "filename": filename,
                "format": source.get("format", source.get("type", "unknown")),
                "chunks": source["chunks"],
                "characters": source["characters"],
                "status": "completed"
            }
        except Exception as e:
            # Clean up source on failure
            await source_store.delete(notebook_id, source["id"])
            raise e

    async def _extract_text(self, content: bytes, filename: str) -> str:
        """Extract text from document with universal fallback for unknown types."""
        file_type = self._get_file_type(filename, content)

        if file_type == "pdf":
            return await self._extract_from_pdf(content)
        elif file_type == "docx":
            return await self._extract_from_docx(content)
        elif file_type == "doc":
            return await self._extract_from_doc_legacy(content)
        elif file_type in ["xlsx", "xls"]:
            return await self._extract_from_excel(content, file_type)
        elif file_type == "csv":
            return await self._extract_from_csv(content)
        elif file_type == "pptx":
            return await self._extract_from_pptx(content)
        elif file_type == "ppt":
            return await self._extract_from_ppt_legacy(content)
        elif file_type in ["mp3", "wav", "m4a", "ogg", "flac", "aac", "wma"]:
            return await self._extract_from_audio(content, filename)
        elif file_type in ["mp4", "mov", "avi", "mkv", "webm", "wmv", "flv"]:
            return await self._extract_from_video(content, filename)
        elif file_type == "epub":
            return await self._extract_from_epub(content)
        elif file_type == "ipynb":
            return await self._extract_from_jupyter(content)
        elif file_type == "odt":
            return await self._extract_from_odt(content)
        elif file_type == "rtf":
            return await self._extract_from_rtf(content)
        elif file_type in ["png", "jpg", "jpeg", "tiff", "bmp", "gif"]:
            return await self._extract_from_image_ocr(content, filename)
        elif file_type in ["html", "htm"]:
            return await self._extract_from_html(content)
        elif file_type in ["txt", "md", "markdown", "json", "xml", "py", "js", "ts", "css", "yaml", "yml", "tex", "bib"]:
            # Text-based files (including LaTeX and BibTeX)
            try:
                return content.decode('utf-8')
            except UnicodeDecodeError:
                return content.decode('latin-1')
        else:
            # Universal fallback: try multiple extraction strategies
            return await self._extract_with_fallback(content, filename, file_type)

    async def _extract_from_pdf(self, content: bytes) -> str:
        """Extract text from PDF with page markers and table handling.
        
        Uses PyMuPDF's table detection when available for better structured data extraction.
        """
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=content, filetype="pdf")
            text_parts = []
            
            for page_num, page in enumerate(doc, 1):
                page_content = [f"=== Page {page_num} ==="]
                
                # Try to extract tables first (PyMuPDF 1.23.0+)
                try:
                    tables = page.find_tables()
                    if tables and tables.tables:
                        for table in tables.tables:
                            table_text = self._extract_pdf_table(table)
                            if table_text:
                                page_content.append(table_text)
                except (AttributeError, Exception):
                    pass  # Table detection not available or failed
                
                # Extract regular text
                page_text = page.get_text().strip()
                if page_text:
                    page_content.append(page_text)
                
                text_parts.append("\n".join(page_content))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to process PDF: {str(e)}")
    
    def _extract_pdf_table(self, table) -> str:
        """Extract table from PDF and convert to searchable text."""
        try:
            import pandas as pd
            
            # PyMuPDF table.extract() returns list of lists
            data = table.extract()
            if not data or len(data) < 2:
                return ""
            
            # First row as headers
            df = pd.DataFrame(data[1:], columns=data[0])
            
            # Clean up empty columns
            df = df.loc[:, (df != '').any(axis=0)]
            
            sentences = self._dataframe_to_sentences(df, "PDF Table")
            return sentences if sentences else ""
        except Exception:
            return ""

    async def _extract_from_docx(self, content: bytes) -> str:
        """Extract text from DOCX including tables, headers, and footers.
        
        Tables are converted to natural language for better semantic search.
        """
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            text_parts = []
            
            # Extract headers
            for section in doc.sections:
                header = section.header
                if header:
                    for para in header.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)
            
            # Extract main content - paragraphs and tables in document order
            for element in doc.element.body:
                if element.tag.endswith('p'):  # Paragraph
                    for para in doc.paragraphs:
                        if para._element == element and para.text.strip():
                            text_parts.append(para.text)
                            break
                elif element.tag.endswith('tbl'):  # Table
                    for table in doc.tables:
                        if table._element == element:
                            table_text = self._extract_docx_table(table)
                            if table_text:
                                text_parts.append(table_text)
                            break
            
            # Fallback: if element iteration didn't work, use simple extraction
            if not text_parts:
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
                for table in doc.tables:
                    table_text = self._extract_docx_table(table)
                    if table_text:
                        text_parts.append(table_text)
            
            # Extract footers
            for section in doc.sections:
                footer = section.footer
                if footer:
                    for para in footer.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to process DOCX: {str(e)}")
    
    async def _extract_from_doc_legacy(self, content: bytes) -> str:
        """Extract text from legacy .doc files (Word 97-2003).
        
        Uses olefile to read OLE compound documents and extract text streams.
        """
        try:
            import olefile
            import re
            
            ole = olefile.OleFileIO(io.BytesIO(content))
            
            # Word documents store text in 'WordDocument' stream
            # The text is encoded in a complex format, but we can extract readable parts
            text_parts = []
            
            # Try to get the main document stream
            if ole.exists('WordDocument'):
                word_stream = ole.openstream('WordDocument').read()
                # Extract ASCII text (basic but works for most content)
                text = word_stream.decode('latin-1', errors='ignore')
                # Clean up binary garbage
                text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', ' ', text)
                text = re.sub(r'\s+', ' ', text)
                # Extract readable sentences (at least 3 words)
                sentences = re.findall(r'[A-Za-z][A-Za-z0-9\s,.\'-]{10,}[.!?]', text)
                if sentences:
                    text_parts.extend(sentences)
            
            # Also try to extract from other streams that might contain text
            for stream_path in ole.listdir():
                stream_name = '/'.join(stream_path)
                if 'text' in stream_name.lower() or 'content' in stream_name.lower():
                    try:
                        stream_data = ole.openstream(stream_path).read()
                        stream_text = stream_data.decode('utf-8', errors='ignore')
                        if stream_text.strip():
                            text_parts.append(stream_text)
                    except:
                        pass
            
            ole.close()
            
            if text_parts:
                return "\n\n".join(text_parts)
            else:
                raise ValueError("No readable text found in DOC file")
                
        except ImportError:
            # Fallback: try to extract any readable text
            import re
            text = content.decode('latin-1', errors='ignore')
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', ' ', text)
            text = re.sub(r'\s+', ' ', text)
            sentences = re.findall(r'[A-Za-z][A-Za-z0-9\s,.\'-]{10,}[.!?]', text)
            if sentences:
                return "\n\n".join(sentences)
            raise ValueError("Legacy DOC processing requires olefile. Install with: pip install olefile")
        except Exception as e:
            raise ValueError(f"Failed to process legacy DOC: {str(e)}")
    
    async def _extract_from_ppt_legacy(self, content: bytes) -> str:
        """Extract text from legacy .ppt files (PowerPoint 97-2003).
        
        Uses olefile to read OLE compound documents and extract text.
        """
        try:
            import olefile
            import re
            
            ole = olefile.OleFileIO(io.BytesIO(content))
            text_parts = []
            
            # PowerPoint stores text in various streams
            for stream_path in ole.listdir():
                stream_name = '/'.join(stream_path)
                try:
                    stream_data = ole.openstream(stream_path).read()
                    # Try to decode as text
                    try:
                        text = stream_data.decode('utf-16-le', errors='ignore')
                    except:
                        text = stream_data.decode('latin-1', errors='ignore')
                    
                    # Clean up and extract readable content
                    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', ' ', text)
                    text = re.sub(r'\s+', ' ', text)
                    
                    # Look for readable sentences
                    sentences = re.findall(r'[A-Za-z][A-Za-z0-9\s,.\'-]{10,}[.!?]', text)
                    text_parts.extend(sentences)
                except:
                    pass
            
            ole.close()
            
            if text_parts:
                # Deduplicate while preserving order
                seen = set()
                unique_parts = []
                for part in text_parts:
                    if part not in seen:
                        seen.add(part)
                        unique_parts.append(part)
                return "\n\n".join(unique_parts)
            else:
                raise ValueError("No readable text found in PPT file")
                
        except ImportError:
            # Fallback: try to extract any readable text
            import re
            text = content.decode('latin-1', errors='ignore')
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', ' ', text)
            sentences = re.findall(r'[A-Za-z][A-Za-z0-9\s,.\'-]{10,}[.!?]', text)
            if sentences:
                return "\n\n".join(sentences)
            raise ValueError("Legacy PPT processing requires olefile. Install with: pip install olefile")
        except Exception as e:
            raise ValueError(f"Failed to process legacy PPT: {str(e)}")
    
    def _extract_docx_table(self, table) -> str:
        """Extract table from DOCX and convert to searchable text."""
        try:
            import pandas as pd
            
            # Extract table data
            data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                data.append(row_data)
            
            if not data:
                return ""
            
            # Try to use first row as headers
            if len(data) > 1:
                df = pd.DataFrame(data[1:], columns=data[0])
            else:
                df = pd.DataFrame(data)
            
            # Convert to sentences for better semantic search
            sentences = self._dataframe_to_sentences(df, "Document Table")
            return sentences if sentences else df.to_string(index=False)
        except Exception:
            # Fallback to simple text extraction
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            return "\n".join(rows)

    async def _extract_from_excel(self, content: bytes, file_type: str) -> str:
        """Extract text from Excel files.
        
        Converts tabular data to natural language sentences for better semantic search.
        Raw tables with column headers create poor embeddings; sentence format works much better.
        """
        try:
            import pandas as pd
            
            # Read Excel file
            if file_type == "xlsx":
                df_dict = pd.read_excel(io.BytesIO(content), sheet_name=None, engine='openpyxl')
            else:  # xls
                df_dict = pd.read_excel(io.BytesIO(content), sheet_name=None, engine='xlrd')
            
            text_parts = []
            for sheet_name, df in df_dict.items():
                text_parts.append(f"=== Sheet: {sheet_name} ===\n")
                
                # Clean up column names - replace "Unnamed: X" with empty string
                df.columns = ['' if 'Unnamed' in str(col) else str(col) for col in df.columns]
                
                # Try to convert to natural language sentences for better semantic search
                # This helps embedding models understand tabular data
                sentences = self._dataframe_to_sentences(df, sheet_name)
                if sentences:
                    text_parts.append(sentences)
                else:
                    # Fallback to table format if sentence conversion fails
                    text_parts.append(df.fillna('').to_string(index=False))
                
                text_parts.append("\n\n")
            
            return "\n".join(text_parts)
        except ImportError:
            raise ValueError("Excel processing requires pandas and openpyxl. Install with: pip install pandas openpyxl xlrd")
        except Exception as e:
            raise ValueError(f"Failed to process Excel file: {str(e)}")
    
    def _dataframe_to_sentences(self, df, sheet_name: str) -> str:
        """Convert a DataFrame to natural language sentences for better semantic embeddings.
        
        Uses multiple strategies:
        1. For pivot-table style data (Q1 FY2025, etc.): Creates "Person - Activity: Q1 FY2025 = 7" sentences
        2. For general tabular data: Creates "Row N: Column1 = Value1, Column2 = Value2" sentences
        3. Fallback: Returns formatted table as text
        
        This ensures ALL data is captured for RAG retrieval.
        """
        import pandas as pd
        import re
        
        sentences = []
        df = df.fillna('')
        
        # Add sheet context
        sentences.append(f"Data from sheet '{sheet_name}':")
        
        # Build column index mapping to handle duplicate column names
        col_to_idx = {}
        for i, col in enumerate(df.columns):
            if col not in col_to_idx:
                col_to_idx[col] = i
        
        def get_cell_value(row_idx, col_idx):
            """Get a scalar value from a cell by position"""
            try:
                val = df.iloc[row_idx, col_idx]
                if hasattr(val, 'iloc'):
                    val = val.iloc[0] if len(val) > 0 else ''
                return val
            except Exception:
                return ''
        
        # Clean column names - replace "Unnamed: X" with position-based names
        clean_col_names = []
        for i, col in enumerate(df.columns):
            col_name = str(col).strip()
            if not col_name or 'Unnamed' in col_name:
                clean_col_names.append(f"Column{i+1}")
            else:
                clean_col_names.append(col_name)
        
        # Check if this looks like pivot-table style data (has Q/FY columns)
        value_pattern = re.compile(r'(Q\s*\d|FY\s*\d|Total|Count|Sum|Amount)', re.IGNORECASE)
        has_pivot_columns = any(value_pattern.search(str(col)) for col in df.columns)
        
        if has_pivot_columns:
            # Pivot-table extraction - identify value columns vs label columns
            label_cols_idx = []
            value_cols_idx = []
            
            for i, col in enumerate(df.columns):
                col_name = str(col).strip()
                if value_pattern.search(col_name):
                    value_cols_idx.append(i)
                else:
                    # Check if this column contains mostly text (label) or numbers (value)
                    # Include ALL non-value columns as potential labels, even "Unnamed" ones
                    try:
                        col_data = df.iloc[:, i].dropna()
                        if len(col_data) > 0:
                            # Check if column has text values (not just numbers)
                            text_count = sum(1 for v in col_data.head(10) 
                                           if isinstance(v, str) and v.strip() and 
                                           not v.replace('.','').replace('-','').isdigit())
                            if text_count > 0:
                                label_cols_idx.append(i)
                            elif df.iloc[:, i].dtype in ['int64', 'float64', 'int32', 'float32']:
                                value_cols_idx.append(i)
                    except Exception:
                        pass
            
            # Track last non-empty values for hierarchical labels
            last_labels = [''] * len(label_cols_idx)
            
            for row_idx in range(len(df)):
                # Update labels - carry forward for merged cells
                current_labels = []
                for j, col_idx in enumerate(label_cols_idx):
                    val = str(get_cell_value(row_idx, col_idx)).strip()
                    if val and val not in ('0', '0.0', 'nan', 'None', ''):
                        last_labels[j] = val
                    if last_labels[j]:
                        current_labels.append(last_labels[j])
                
                if not current_labels:
                    continue
                
                row_subject = " - ".join(current_labels)
                
                # Extract values
                for col_idx in value_cols_idx:
                    col_name = clean_col_names[col_idx]
                    val = get_cell_value(row_idx, col_idx)
                    try:
                        if pd.isna(val) or val == '' or val == 0:
                            continue
                        num_val = float(val) if not isinstance(val, (int, float)) else val
                        if num_val == 0:
                            continue
                        formatted_val = int(num_val) if float(num_val).is_integer() else num_val
                        sentences.append(f"{row_subject}: {col_name} = {formatted_val}")
                    except (ValueError, TypeError):
                        pass
        
        # ALWAYS also add row-by-row extraction for complete coverage
        # This ensures no data is missed even if pivot extraction fails
        sentences.append("")  # Blank line separator
        sentences.append("Complete row data:")
        
        for row_idx in range(len(df)):
            row_parts = []
            for col_idx in range(len(df.columns)):
                val = get_cell_value(row_idx, col_idx)
                val_str = str(val).strip()
                if val_str and val_str not in ('nan', 'None', ''):
                    col_name = clean_col_names[col_idx]
                    row_parts.append(f"{col_name}: {val_str}")
            
            if row_parts:
                sentences.append(f"Row {row_idx + 1}: " + ", ".join(row_parts))
        
        return "\n".join(sentences) if len(sentences) > 2 else ""

    async def _extract_from_csv(self, content: bytes) -> str:
        """Extract text from CSV files.
        
        Converts tabular data to natural language sentences for better semantic search.
        """
        try:
            import pandas as pd
            
            # Try to detect encoding
            try:
                text = content.decode('utf-8')
            except UnicodeDecodeError:
                text = content.decode('latin-1')
            
            df = pd.read_csv(io.StringIO(text))
            
            # Convert to sentences for better semantic search
            sentences = self._dataframe_to_sentences(df, "CSV Data")
            if sentences:
                return sentences
            return df.fillna('').to_string(index=False)
        except ImportError:
            # Fallback without pandas
            try:
                return content.decode('utf-8')
            except UnicodeDecodeError:
                return content.decode('latin-1')
        except Exception as e:
            raise ValueError(f"Failed to process CSV file: {str(e)}")

    async def _extract_from_pptx(self, content: bytes) -> str:
        """Extract text from PowerPoint files including tables and speaker notes.
        
        Tables are converted to natural language for better semantic search.
        """
        try:
            from pptx import Presentation
            from pptx.shapes.table import Table
            
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"=== Slide {slide_num} ==="]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        slide_content.append(shape.text)
                    
                    # Extract tables
                    if shape.has_table:
                        table_text = self._extract_pptx_table(shape.table)
                        if table_text:
                            slide_content.append(table_text)
                
                # Extract speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_content.append(f"Speaker Notes: {notes}")
                
                text_parts.extend(slide_content)
                text_parts.append("")  # Blank line between slides
            
            return "\n".join(text_parts)
        except ImportError:
            raise ValueError("PowerPoint processing requires python-pptx. Install with: pip install python-pptx")
        except Exception as e:
            raise ValueError(f"Failed to process PowerPoint file: {str(e)}")
    
    def _extract_pptx_table(self, table) -> str:
        """Extract table from PPTX and convert to searchable text."""
        try:
            import pandas as pd
            
            data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                data.append(row_data)
            
            if not data:
                return ""
            
            if len(data) > 1:
                df = pd.DataFrame(data[1:], columns=data[0])
            else:
                df = pd.DataFrame(data)
            
            sentences = self._dataframe_to_sentences(df, "Slide Table")
            return sentences if sentences else df.to_string(index=False)
        except Exception:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            return "\n".join(rows)

    async def _extract_from_audio(self, content: bytes, filename: str) -> str:
        """Extract text from audio files using speech-to-text"""
        import tempfile
        import os
        
        try:
            import whisper
            
            # Save to temp file (whisper needs file path)
            with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            
            try:
                # Load openai-whisper model (uses 'base' for speed, can use 'small' or 'medium' for accuracy)
                model = whisper.load_model("base")
                result = model.transcribe(tmp_path)
                return result["text"]
            finally:
                # Clean up temp file
                os.unlink(tmp_path)
                
        except ImportError:
            raise ValueError("Audio transcription requires openai-whisper. Install with: pip install openai-whisper")
        except Exception as e:
            raise ValueError(f"Failed to transcribe audio: {str(e)}")

    async def _extract_from_video(self, content: bytes, filename: str) -> str:
        """Extract text from video files by extracting audio and transcribing"""
        import tempfile
        import subprocess
        import os
        
        try:
            import whisper
            
            # Save video to temp file
            with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as tmp_video:
                tmp_video.write(content)
                video_path = tmp_video.name
            
            # Extract audio using ffmpeg
            audio_path = video_path + ".wav"
            
            try:
                subprocess.run(
                    ["ffmpeg", "-i", video_path, "-vn", "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1", "-y", audio_path],
                    check=True,
                    capture_output=True,
                    timeout=300  # 5 min timeout for long videos
                )
                
                # Transcribe the extracted audio
                model = whisper.load_model("base")
                result = model.transcribe(audio_path)
                return result["text"]
                
            finally:
                # Clean up temp files
                if os.path.exists(video_path):
                    os.unlink(video_path)
                if os.path.exists(audio_path):
                    os.unlink(audio_path)
                    
        except ImportError:
            raise ValueError("Video transcription requires openai-whisper and ffmpeg. Install with: pip install openai-whisper")
        except subprocess.CalledProcessError:
            raise ValueError("Video processing requires ffmpeg. Install with: brew install ffmpeg")
        except Exception as e:
            raise ValueError(f"Failed to transcribe video: {str(e)}")

    async def _extract_from_epub(self, content: bytes) -> str:
        """Extract text from EPUB e-books."""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            
            book = epub.read_epub(io.BytesIO(content))
            text_parts = []
            
            # Extract metadata
            title = book.get_metadata('DC', 'title')
            if title:
                text_parts.append(f"Title: {title[0][0]}")
            
            author = book.get_metadata('DC', 'creator')
            if author:
                text_parts.append(f"Author: {author[0][0]}")
            
            text_parts.append("")
            
            # Extract content from each chapter
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    chapter_text = soup.get_text(separator='\n', strip=True)
                    if chapter_text:
                        text_parts.append(chapter_text)
            
            return "\n\n".join(text_parts)
        except ImportError:
            raise ValueError("EPUB processing requires ebooklib. Install with: pip install ebooklib")
        except Exception as e:
            raise ValueError(f"Failed to process EPUB: {str(e)}")
    
    async def _extract_from_jupyter(self, content: bytes) -> str:
        """Extract text from Jupyter notebooks (.ipynb)."""
        try:
            import nbformat
            import json
            
            notebook = nbformat.reads(content.decode('utf-8'), as_version=4)
            text_parts = []
            
            for cell_num, cell in enumerate(notebook.cells, 1):
                if cell.cell_type == 'markdown':
                    text_parts.append(f"[Markdown Cell {cell_num}]")
                    text_parts.append(cell.source)
                elif cell.cell_type == 'code':
                    text_parts.append(f"[Code Cell {cell_num}]")
                    text_parts.append(f"```python\n{cell.source}\n```")
                    # Include outputs if they contain text
                    for output in cell.get('outputs', []):
                        if output.get('output_type') == 'stream':
                            text_parts.append(f"Output: {output.get('text', '')}")
                        elif output.get('output_type') == 'execute_result':
                            data = output.get('data', {})
                            if 'text/plain' in data:
                                text_parts.append(f"Result: {data['text/plain']}")
                text_parts.append("")
            
            return "\n".join(text_parts)
        except ImportError:
            raise ValueError("Jupyter notebook processing requires nbformat. Install with: pip install nbformat")
        except Exception as e:
            raise ValueError(f"Failed to process Jupyter notebook: {str(e)}")
    
    async def _extract_from_odt(self, content: bytes) -> str:
        """Extract text from OpenDocument Text files (.odt) including tables."""
        try:
            from odf import text as odf_text
            from odf import table as odf_table
            from odf.opendocument import load
            
            doc = load(io.BytesIO(content))
            text_parts = []
            
            # Extract paragraphs
            for para in doc.getElementsByType(odf_text.P):
                para_text = self._extract_odt_text_content(para)
                if para_text.strip():
                    text_parts.append(para_text)
            
            # Extract tables
            for table in doc.getElementsByType(odf_table.Table):
                table_text = self._extract_odt_table(table)
                if table_text:
                    text_parts.append(table_text)
            
            return "\n\n".join(text_parts)
        except ImportError:
            raise ValueError("ODT processing requires odfpy. Install with: pip install odfpy")
        except Exception as e:
            raise ValueError(f"Failed to process ODT: {str(e)}")
    
    def _extract_odt_text_content(self, element) -> str:
        """Recursively extract text from ODT element."""
        text = ""
        for node in element.childNodes:
            if node.nodeType == node.TEXT_NODE:
                text += node.data
            elif hasattr(node, 'childNodes'):
                text += self._extract_odt_text_content(node)
        return text
    
    def _extract_odt_table(self, table) -> str:
        """Extract table from ODT and convert to searchable text."""
        try:
            from odf import table as odf_table
            import pandas as pd
            
            rows_data = []
            for row in table.getElementsByType(odf_table.TableRow):
                row_data = []
                for cell in row.getElementsByType(odf_table.TableCell):
                    cell_text = self._extract_odt_text_content(cell).strip()
                    row_data.append(cell_text)
                if row_data:
                    rows_data.append(row_data)
            
            if not rows_data:
                return ""
            
            # Try to use first row as headers
            if len(rows_data) > 1:
                df = pd.DataFrame(rows_data[1:], columns=rows_data[0])
            else:
                df = pd.DataFrame(rows_data)
            
            sentences = self._dataframe_to_sentences(df, "ODT Table")
            return sentences if sentences else df.to_string(index=False)
        except Exception:
            # Fallback to simple extraction
            rows = []
            for row in rows_data:
                rows.append(" | ".join(row))
            return "\n".join(rows)
    
    async def _extract_from_rtf(self, content: bytes) -> str:
        """Extract text from RTF files using striprtf library."""
        try:
            from striprtf.striprtf import rtf_to_text
            
            # Decode RTF content
            rtf_content = content.decode('utf-8', errors='ignore')
            
            # Use striprtf for proper RTF parsing
            text = rtf_to_text(rtf_content)
            
            # Clean up any remaining artifacts
            import re
            # Remove any remaining control sequences that slipped through
            text = re.sub(r'\\[a-z]+\d*\s?', ' ', text)
            # Normalize whitespace
            text = re.sub(r'\s+', ' ', text)
            # Remove isolated backslashes
            text = re.sub(r'\\\s', ' ', text)
            
            return text.strip()
        except ImportError:
            # Fallback to basic extraction if striprtf not installed
            print("[WARN] striprtf not installed, using basic RTF extraction")
            text = content.decode('utf-8', errors='ignore')
            import re
            text = re.sub(r'\\[a-z]+\d*\s?', ' ', text)
            text = re.sub(r'[{}]', '', text)
            text = re.sub(r'\\\'[0-9a-f]{2}', '', text)
            text = re.sub(r'\s+', ' ', text)
            return text.strip()
        except Exception as e:
            raise ValueError(f"Failed to process RTF: {str(e)}")
    
    async def _extract_from_html(self, content: bytes) -> str:
        """Extract text from HTML files with proper tag stripping."""
        try:
            from bs4 import BeautifulSoup
            
            # Decode content
            try:
                html_content = content.decode('utf-8')
            except UnicodeDecodeError:
                html_content = content.decode('latin-1')
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style", "meta", "link", "noscript"]):
                script.decompose()
            
            # Extract title if present
            text_parts = []
            title = soup.find('title')
            if title and title.string:
                text_parts.append(f"Title: {title.string.strip()}")
            
            # Get text with proper spacing
            text = soup.get_text(separator='\n', strip=True)
            
            # Clean up excessive whitespace
            import re
            text = re.sub(r'\n\s*\n', '\n\n', text)
            text_parts.append(text)
            
            return "\n\n".join(text_parts)
        except ImportError:
            # Fallback without BeautifulSoup - basic tag stripping
            import re
            try:
                text = content.decode('utf-8')
            except UnicodeDecodeError:
                text = content.decode('latin-1')
            text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<[^>]+>', ' ', text)
            text = re.sub(r'\s+', ' ', text)
            return text.strip()
        except Exception as e:
            raise ValueError(f"Failed to process HTML: {str(e)}")
    
    async def _extract_from_image_ocr(self, content: bytes, filename: str) -> str:
        """Extract text from images using OCR (Tesseract)."""
        try:
            import pytesseract
            from PIL import Image
            
            image = Image.open(io.BytesIO(content))
            
            # Perform OCR
            text = pytesseract.image_to_string(image)
            
            if not text.strip():
                return f"[Image: {filename} - No text detected via OCR]"
            
            return f"[OCR from {filename}]\n{text}"
        except ImportError:
            raise ValueError("Image OCR requires pytesseract and Pillow. Install with: pip install pytesseract Pillow. Also install Tesseract: brew install tesseract")
        except Exception as e:
            raise ValueError(f"Failed to OCR image: {str(e)}")

    async def _extract_with_fallback(self, content: bytes, filename: str, file_type: str) -> str:
        """Universal fallback extractor - tries multiple strategies for unknown file types.
        
        This ensures we can extract text from files even when:
        - The file type is unknown or unsupported
        - The file has a wrong extension
        - The file is a new format we haven't explicitly handled
        """
        print(f"[DocProcessor] Using fallback extraction for {filename} (type: {file_type})")
        
        extraction_attempts = []
        
        # Strategy 1: Try as plain text (UTF-8, then Latin-1)
        try:
            text = content.decode('utf-8')
            if text.strip() and len(text.strip()) > 20:
                print(f"[DocProcessor] Fallback: decoded as UTF-8 text")
                return text
        except UnicodeDecodeError:
            try:
                text = content.decode('latin-1')
                if text.strip() and len(text.strip()) > 20:
                    print(f"[DocProcessor] Fallback: decoded as Latin-1 text")
                    return text
            except:
                pass
        extraction_attempts.append("text decode")
        
        # Strategy 2: Try as PDF (common for mislabeled files)
        try:
            text = await self._extract_from_pdf(content)
            if text and len(text.strip()) > 50:
                print(f"[DocProcessor] Fallback: extracted as PDF")
                return text
        except:
            pass
        extraction_attempts.append("PDF")
        
        # Strategy 3: Try as DOCX
        try:
            text = await self._extract_from_docx(content)
            if text and len(text.strip()) > 50:
                print(f"[DocProcessor] Fallback: extracted as DOCX")
                return text
        except:
            pass
        extraction_attempts.append("DOCX")
        
        # Strategy 4: Try as Excel
        try:
            text = await self._extract_from_excel(content, 'xlsx')
            if text and len(text.strip()) > 50:
                print(f"[DocProcessor] Fallback: extracted as Excel")
                return text
        except:
            pass
        extraction_attempts.append("Excel")
        
        # Strategy 5: Try OCR as last resort (treat as image)
        try:
            text = await self._extract_from_image_ocr(content, filename)
            if text and len(text.strip()) > 20:
                print(f"[DocProcessor] Fallback: extracted via OCR")
                return text
        except:
            pass
        extraction_attempts.append("OCR")
        
        # All strategies failed
        raise ValueError(f"Could not extract text from {filename} (type: {file_type}). Tried: {', '.join(extraction_attempts)}")

    def _get_file_type(self, filename: str, content: bytes = None) -> str:
        """Get file type from filename, with magic byte fallback for unknown/missing extensions.
        
        This ensures files with wrong or missing extensions are still processed correctly.
        """
        ext = Path(filename).suffix.lower()
        ext = ext[1:] if ext else ""
        
        # If we have a known extension, use it
        known_extensions = {
            'pdf', 'docx', 'doc', 'xlsx', 'xls', 'csv', 'pptx', 'ppt',
            'txt', 'md', 'markdown', 'json', 'xml', 'html', 'htm',
            'py', 'js', 'ts', 'css', 'yaml', 'yml', 'tex', 'bib',
            'epub', 'ipynb', 'odt', 'rtf',
            'mp3', 'wav', 'm4a', 'ogg', 'flac', 'aac', 'wma',
            'mp4', 'mov', 'avi', 'mkv', 'webm', 'wmv', 'flv',
            'png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif'
        }
        
        if ext in known_extensions:
            return ext
        
        # Fallback: detect by magic bytes (file signature)
        if content and len(content) >= 8:
            detected = self._detect_by_magic_bytes(content)
            if detected:
                print(f"[DocProcessor] Detected file type by magic bytes: {detected} (filename had: {ext or 'no extension'})")
                return detected
        
        # Last resort: return extension or unknown
        return ext if ext else "unknown"
    
    def _detect_by_magic_bytes(self, content: bytes) -> str:
        """Detect file type by magic bytes (file signature).
        
        This handles files with wrong or missing extensions.
        """
        # Common file signatures (magic bytes)
        signatures = [
            # Documents
            (b'%PDF', 'pdf'),
            (b'PK\x03\x04', '_zip'),  # ZIP-based (docx, xlsx, pptx, epub, odt)
            (b'\xd0\xcf\x11\xe0', '_ole'),  # OLE (doc, xls, ppt)
            (b'{\\rtf', 'rtf'),
            
            # Images
            (b'\x89PNG\r\n\x1a\n', 'png'),
            (b'\xff\xd8\xff', 'jpg'),
            (b'GIF87a', 'gif'),
            (b'GIF89a', 'gif'),
            (b'BM', 'bmp'),
            (b'II*\x00', 'tiff'),  # Little-endian TIFF
            (b'MM\x00*', 'tiff'),  # Big-endian TIFF
            
            # Audio
            (b'ID3', 'mp3'),
            (b'\xff\xfb', 'mp3'),
            (b'\xff\xfa', 'mp3'),
            (b'RIFF', '_riff'),  # WAV or AVI
            (b'fLaC', 'flac'),
            (b'OggS', 'ogg'),
            
            # Video
            (b'\x00\x00\x00\x1cftyp', 'mp4'),
            (b'\x00\x00\x00\x20ftyp', 'mp4'),
            (b'\x1aE\xdf\xa3', 'mkv'),
        ]
        
        for sig, file_type in signatures:
            if content.startswith(sig):
                # Handle compound formats
                if file_type == '_zip':
                    return self._detect_zip_subtype(content)
                elif file_type == '_ole':
                    return self._detect_ole_subtype(content)
                elif file_type == '_riff':
                    return self._detect_riff_subtype(content)
                return file_type
        
        # Check if it's valid UTF-8 text
        try:
            content[:1000].decode('utf-8')
            # Check for JSON
            if content.strip().startswith(b'{') or content.strip().startswith(b'['):
                return 'json'
            # Check for XML/HTML
            if content.strip().startswith(b'<?xml') or content.strip().startswith(b'<'):
                if b'<html' in content[:500].lower():
                    return 'html'
                return 'xml'
            return 'txt'
        except UnicodeDecodeError:
            pass
        
        return None
    
    def _detect_zip_subtype(self, content: bytes) -> str:
        """Detect specific type of ZIP-based file (docx, xlsx, pptx, epub, odt)."""
        try:
            import zipfile
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                names = zf.namelist()
                
                # Check for Office Open XML
                if '[Content_Types].xml' in names:
                    if any('word/' in n for n in names):
                        return 'docx'
                    if any('xl/' in n for n in names):
                        return 'xlsx'
                    if any('ppt/' in n for n in names):
                        return 'pptx'
                
                # Check for EPUB
                if 'META-INF/container.xml' in names:
                    return 'epub'
                
                # Check for ODT
                if 'content.xml' in names and 'mimetype' in names:
                    return 'odt'
                
                # Check for Jupyter notebook
                if any(n.endswith('.ipynb') for n in names):
                    return 'ipynb'
        except:
            pass
        return 'zip'
    
    def _detect_ole_subtype(self, content: bytes) -> str:
        """Detect specific type of OLE file (doc, xls, ppt)."""
        # OLE files are complex; default to doc as most common
        # Could use olefile library for precise detection
        return 'doc'
    
    def _detect_riff_subtype(self, content: bytes) -> str:
        """Detect RIFF subtype (WAV vs AVI)."""
        if len(content) >= 12:
            riff_type = content[8:12]
            if riff_type == b'WAVE':
                return 'wav'
            if riff_type == b'AVI ':
                return 'avi'
        return 'wav'  # Default to WAV

document_processor = DocumentProcessor()
